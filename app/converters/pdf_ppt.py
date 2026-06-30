import os
from pathlib import Path
from io import BytesIO

def convert_pdf_to_ppt(input_path, output_path):
    """Convert PDF to PowerPoint.
    Tries PyMuPDF first, falls back to pure Python.
    """
    try:
        return _convert_ppt_with_pymupdf(input_path, output_path)
    except Exception as e:
        print(f"[PyMuPDF PPT failed: {e}], trying pure Python fallback...")

    try:
        return _convert_ppt_with_pypdf(input_path, output_path)
    except Exception as e:
        raise RuntimeError(f"Both PyMuPDF and pure Python failed: {e}")

def _convert_ppt_with_pymupdf(input_path, output_path):
    """Fast PPT conversion using PyMuPDF."""
    import fitz
    from pptx import Presentation
    from pptx.util import Inches
    from PIL import Image as PILImage

    doc = None
    try:
        doc = fitz.open(str(input_path))
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        blank_slide_layout = prs.slide_layouts[6]

        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            slide = prs.slides.add_slide(blank_slide_layout)

            pix = page.get_pixmap()
            img = PILImage.frombytes("RGB", [pix.width, pix.height], pix.samples)

            img_buffer = BytesIO()
            img.save(img_buffer, format="PNG")
            img_buffer.seek(0)

            left = Inches(0)
            top = Inches(0)
            width = prs.slide_width
            height = prs.slide_height

            slide.shapes.add_picture(img_buffer, left, top, width=width, height=height)

            text = page.get_text()
            if text and text.strip():
                txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
                tf = txBox.text_frame
                tf.text = text[:500]

        prs.save(str(output_path))
    finally:
        if doc is not None:
            try:
                doc.close()
            except Exception:
                pass

def _convert_ppt_with_pypdf(input_path, output_path):
    """Pure Python PPT fallback."""
    from pypdf import PdfReader
    from pptx import Presentation
    from pptx.util import Inches
    from PIL import Image as PILImage, ImageDraw, ImageFont

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_slide_layout = prs.slide_layouts[6]
    reader = PdfReader(str(input_path))

    for page_num, page in enumerate(reader.pages):
        slide = prs.slides.add_slide(blank_slide_layout)

        try:
            width = float(page.mediabox.width)
            height = float(page.mediabox.height)
        except Exception:
            width, height = 612.0, 792.0

        scale = 150 / 72.0
        img_width = max(int(width * scale), 100)
        img_height = max(int(height * scale), 100)

        img = PILImage.new('RGB', (img_width, img_height), (255, 255, 255))
        draw = ImageDraw.Draw(img)

        try:
            font = ImageFont.truetype("arial.ttf", 14)
        except Exception:
            font = ImageFont.load_default()

        text = None
        try:
            text = page.extract_text()
        except Exception:
            try:
                text = page.extract_text(extraction_mode="layout")
            except Exception:
                text = None

        if text:
            y = 10
            for line in text.split('\n')[:80]:
                line = line.strip()
                if line:
                    if len(line) > 120:
                        line = line[:120] + "..."
                    draw.text((10, y), line, fill=(0, 0, 0), font=font)
                    y += 18
                    if y > img_height - 20:
                        break
        else:
            draw.text((10, 10), f"Page {page_num + 1}", fill=(100, 100, 100), font=font)

        img_buffer = BytesIO()
        img.save(img_buffer, format="PNG")
        img_buffer.seek(0)

        slide.shapes.add_picture(img_buffer, Inches(0), Inches(0), 
                                  width=prs.slide_width, height=prs.slide_height)

        if text and text.strip():
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
            tf = txBox.text_frame
            tf.text = text[:500]

    prs.save(str(output_path))
