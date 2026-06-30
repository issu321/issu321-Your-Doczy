import os
from pathlib import Path
from io import BytesIO

def convert_ppt_to_pdf(input_path, output_path):
    """Convert PowerPoint to PDF."""
    # Try LibreOffice first
    try:
        import subprocess
        cmd = [
            'libreoffice', '--headless', '--convert-to', 'pdf',
            '--outdir', str(Path(output_path).parent),
            str(input_path)
        ]
        subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, timeout=120)
        lo_output = Path(output_path).parent / (Path(input_path).stem + '.pdf')
        if lo_output.exists() and str(lo_output) != str(output_path):
            lo_output.rename(output_path)
        return
    except Exception:
        pass

    # Fallback: python-pptx + reportlab
    _convert_ppt_to_pdf_fallback(input_path, output_path)

def _convert_ppt_to_pdf_fallback(input_path, output_path):
    from pptx import Presentation
    from reportlab.lib.pagesizes import letter, landscape
    from reportlab.pdfgen import canvas
    from reportlab.lib.units import inch
    from PIL import Image as PILImage

    prs = Presentation(input_path)
    c = canvas.Canvas(str(output_path), pagesize=landscape(letter))
    width, height = landscape(letter)

    for slide_num, slide in enumerate(prs.slides):
        # Draw background
        c.setFillColorRGB(1, 1, 1)
        c.rect(0, 0, width, height, fill=1, stroke=0)

        # Extract shapes
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text = shape.text.strip()
                left = shape.left.inches * inch if hasattr(shape.left, 'inches') else shape.left * 0.01 * width
                top = height - (shape.top.inches * inch if hasattr(shape.top, 'inches') else shape.top * 0.01 * height) - 0.5*inch

                c.setFont("Helvetica", 14)
                c.setFillColorRGB(0, 0, 0)
                c.drawString(left, top, text[:200])  # Limit text length

            # Handle images
            if shape.shape_type == 13:  # Picture
                try:
                    image = shape.image
                    image_bytes = image.blob
                    img = PILImage.open(BytesIO(image_bytes))
                    img_path = f"/tmp/ppt_img_{slide_num}.png"
                    img.save(img_path)

                    left = shape.left.inches * inch if hasattr(shape.left, 'inches') else shape.left * 0.01 * width
                    top = height - (shape.top.inches * inch if hasattr(shape.top, 'inches') else shape.top * 0.01 * height) - (shape.height.inches * inch if hasattr(shape.height, 'inches') else shape.height * 0.01 * height)
                    img_width = shape.width.inches * inch if hasattr(shape.width, 'inches') else shape.width * 0.01 * width
                    img_height = shape.height.inches * inch if hasattr(shape.height, 'inches') else shape.height * 0.01 * height

                    c.drawImage(img_path, left, top, width=img_width, height=img_height)
                    if os.path.exists(img_path):
                        os.remove(img_path)
                except Exception:
                    pass

        c.showPage()

    c.save()
